"""QAパネルの資料出力（HTML / Word / PDF）

LLM が生成した Markdown 形式の資料テキストを、配布しやすいファイル形式に
変換する。外部サービスには依存しない（PDF は PyMuPDF Story ＋ Windows の
日本語フォント、Word は python-docx）。
"""

import html as html_lib
import io
import os
import re
from datetime import datetime

logger = __import__("logging").getLogger(__name__)

# PDF 埋め込み用の日本語フォント（見つかったものを使う）
_JP_FONT_CANDIDATES = [
    r"C:\Windows\Fonts\yugothm.ttc",   # 游ゴシック Medium
    r"C:\Windows\Fonts\YuGothM.ttc",
    r"C:\Windows\Fonts\meiryo.ttc",    # メイリオ
    r"C:\Windows\Fonts\msgothic.ttc",  # MSゴシック
]

MIME = {
    "html": "text/html",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf":  "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def safe_filename(title: str, ext: str) -> str:
    """タイトルからファイル名を作る（Windowsで使えない文字を除去）。"""
    name = re.sub(r'[\\/:*?"<>|\r\n]+', "", str(title or "")).strip() or "資料"
    return f"{datetime.now().strftime('%Y%m%d_%H%M')}_{name[:40]}.{ext}"


# ─── Markdown の簡易パース ─────────────────────────────────
# LLM 出力の基本要素（見出し・箇条書き・番号リスト・表・段落・太字）のみ対応。

def _parse_blocks(md: str) -> list:
    """Markdown を [(type, payload), ...] に分解する。
    type: heading(level, text) / bullets([items]) / numbers([items]) /
          table([rows]) / para(text)
    """
    blocks: list = []
    lines = str(md or "").replace("\r\n", "\n").split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        if not line.strip():
            i += 1
            continue
        m = re.match(r"^(#{1,4})\s*(.+)$", line)
        if m:
            blocks.append(("heading", (len(m.group(1)), m.group(2).strip())))
            i += 1
            continue
        if re.match(r"^\s*[-*・]\s+", line):
            items = []
            while i < len(lines) and re.match(r"^\s*[-*・]\s+", lines[i]):
                items.append(re.sub(r"^\s*[-*・]\s+", "", lines[i]).strip())
                i += 1
            blocks.append(("bullets", items))
            continue
        if re.match(r"^\s*\d+[.)、．]\s*", line):
            items = []
            while i < len(lines) and re.match(r"^\s*\d+[.)、．]\s*", lines[i]):
                items.append(re.sub(r"^\s*\d+[.)、．]\s*", "", lines[i]).strip())
                i += 1
            blocks.append(("numbers", items))
            continue
        if line.strip().startswith("|"):
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not all(re.fullmatch(r":?-{2,}:?", c or "-") for c in cells):
                    rows.append(cells)   # 区切り行（|---|---|）は捨てる
                i += 1
            if rows:
                blocks.append(("table", rows))
            continue
        # 連続する通常行は1段落にまとめる
        para = [line.strip()]
        i += 1
        while i < len(lines) and lines[i].strip() \
                and not re.match(r"^(#{1,4}\s|\s*[-*・]\s|\s*\d+[.)、．]\s*|\s*\|)", lines[i]):
            para.append(lines[i].strip())
            i += 1
        blocks.append(("para", "\n".join(para)))
    return blocks


def _inline_html(text: str) -> str:
    """エスケープ後、太字（**）と改行のみHTML化する。"""
    t = html_lib.escape(str(text or ""))
    t = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", t)
    return t.replace("\n", "<br>")


# ─── HTML ─────────────────────────────────────────────

_DOC_CSS = """
  body { font-family: 'Yu Gothic Medium', 'Noto Sans JP', 'メイリオ', sans-serif;
         font-size: 13px; color: #212121; background: #fafafa; margin: 0; padding: 24px; }
  .container { max-width: 860px; margin: 0 auto; background: white;
               padding: 32px 40px; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,.1); }
  h1 { color: #1565C0; border-bottom: 3px solid #1565C0; padding-bottom: 8px; font-size: 20px; }
  h2 { color: #1565C0; font-size: 16px; margin-top: 26px;
       border-left: 4px solid #1565C0; padding-left: 10px; }
  h3 { font-size: 14px; color: #333; margin: 14px 0 6px; }
  p  { line-height: 1.8; }
  li { margin: 3px 0; line-height: 1.7; }
  table { border-collapse: collapse; width: 100%; margin: 10px 0; font-size: 12px; }
  th { background: #E8EAF6; }
  th, td { border: 1px solid #ccc; padding: 6px 10px; text-align: left; }
  .doc-meta { color: #666; font-size: 11px; margin-bottom: 16px; }
  .doc-note { background: #FFF8E1; border-left: 4px solid #FFA000; padding: 10px 14px;
              border-radius: 4px; font-size: 11px; color: #5D4037; margin-top: 28px; }
  @media print { body { padding: 0; background: white; }
                 .container { box-shadow: none; padding: 8px; } }
"""

_DOC_NOTE = ("本資料は「法令・届出施設確認サポートAI」への質問から生成された参考資料です。"
             "内容の正確性は担当部署・所轄機関にご確認のうえご利用ください。")


def _blocks_to_html(blocks: list) -> str:
    out = []
    for kind, payload in blocks:
        if kind == "heading":
            level, text = payload
            out.append(f"<h{min(level, 3)}>{_inline_html(text)}</h{min(level, 3)}>")
        elif kind == "bullets":
            out.append("<ul>" + "".join(f"<li>{_inline_html(x)}</li>" for x in payload) + "</ul>")
        elif kind == "numbers":
            out.append("<ol>" + "".join(f"<li>{_inline_html(x)}</li>" for x in payload) + "</ol>")
        elif kind == "table":
            head, *rest = payload
            out.append(
                "<table><tr>" + "".join(f"<th>{_inline_html(c)}</th>" for c in head) + "</tr>"
                + "".join("<tr>" + "".join(f"<td>{_inline_html(c)}</td>" for c in r) + "</tr>"
                          for r in rest)
                + "</table>"
            )
        else:
            out.append(f"<p>{_inline_html(payload)}</p>")
    return "\n".join(out)


def build_html_file(title: str, md: str) -> bytes:
    body = _blocks_to_html(_parse_blocks(md))
    doc = f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{html_lib.escape(str(title))}</title>
<style>{_DOC_CSS}</style>
</head>
<body>
<div class="container">
<h1>{_inline_html(title)}</h1>
<div class="doc-meta">作成日：{datetime.now().strftime('%Y年%m月%d日')}　（法令・届出施設確認サポートAI）</div>
{body}
<div class="doc-note">⚠️ {_DOC_NOTE}</div>
</div>
</body>
</html>"""
    return doc.encode("utf-8")


# ─── Word (docx) ───────────────────────────────────────

def _docx_add_runs(paragraph, text: str) -> None:
    """**太字** を run に分解して追加する。"""
    for part in re.split(r"(\*\*.+?\*\*)", str(text or "")):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            paragraph.add_run(part[2:-2]).bold = True
        else:
            paragraph.add_run(part)


def build_docx_file(title: str, md: str) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn

    doc = Document()
    # 既定フォントを日本語フォントにする（東アジア文字の指定も必要）
    style = doc.styles["Normal"]
    style.font.name = "游ゴシック"
    style.font.size = Pt(10.5)
    style.element.rPr.rFonts.set(qn("w:eastAsia"), "游ゴシック")

    h = doc.add_heading(str(title), level=0)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x15, 0x65, 0xC0)
    doc.add_paragraph(
        f"作成日：{datetime.now().strftime('%Y年%m月%d日')}　（法令・届出施設確認サポートAI）"
    ).runs[0].font.size = Pt(8)

    for kind, payload in _parse_blocks(md):
        if kind == "heading":
            level, text = payload
            doc.add_heading(re.sub(r"\*\*", "", text), level=min(level, 3))
        elif kind in ("bullets", "numbers"):
            style_name = "List Bullet" if kind == "bullets" else "List Number"
            for item in payload:
                _docx_add_runs(doc.add_paragraph(style=style_name), item)
        elif kind == "table":
            table = doc.add_table(rows=len(payload), cols=max(len(r) for r in payload))
            table.style = "Table Grid"
            for ri, row in enumerate(payload):
                for ci, cell in enumerate(row):
                    p = table.cell(ri, ci).paragraphs[0]
                    if ri == 0:
                        p.add_run(re.sub(r"\*\*", "", cell)).bold = True
                    else:
                        _docx_add_runs(p, cell)
        else:
            _docx_add_runs(doc.add_paragraph(), payload)

    note = doc.add_paragraph()
    note.add_run("⚠️ " + _DOC_NOTE).font.size = Pt(8)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ─── PowerPoint (pptx) ─────────────────────────────────

_PPTX_JP_FONT = "游ゴシック"
_PPTX_MAX_LINES = 8      # 1スライドに載せる本文行数の上限（超えたら「続き」スライド）


def _pptx_set_jp(run, size_pt=None, bold=None, color=None) -> None:
    """run に日本語フォントを設定する（東アジア文字用の typeface も指定）。"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    run.font.name = _PPTX_JP_FONT
    if size_pt:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        from pptx.dml.color import RGBColor
        run.font.color.rgb = RGBColor(*color)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", _PPTX_JP_FONT)


def _pptx_add_runs(paragraph, text: str, size_pt=14) -> None:
    """**太字** を run に分解して段落へ追加する。"""
    for part in re.split(r"(\*\*.+?\*\*)", str(text or "")):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            _pptx_set_jp(run, size_pt, bold=True)
        else:
            run.text = part
            _pptx_set_jp(run, size_pt)


def _split_slides(blocks: list) -> list:
    """ブロック列を (スライドタイトル, [ブロック]) に分割する。
    「##」以下の見出しごとに1枚。最初の見出しより前の内容は「概要」スライド。"""
    slides: list = []
    cur_title, cur_blocks = "概要", []
    for kind, payload in blocks:
        if kind == "heading" and payload[0] <= 2:
            if cur_blocks:
                slides.append((cur_title, cur_blocks))
            cur_title, cur_blocks = re.sub(r"\*\*", "", payload[1]), []
        else:
            cur_blocks.append((kind, payload))
    if cur_blocks:
        slides.append((cur_title, cur_blocks))
    return slides


def build_pptx_file(title: str, md: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    margin = Inches(0.6)
    body_width = prs.slide_width - margin * 2

    def add_slide(slide_title: str, first=False):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(
            margin, Inches(1.0) if not first else Inches(2.6),
            body_width, Inches(1.0))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = slide_title
        _pptx_set_jp(run, 32 if first else 24, bold=True, color=(0x15, 0x65, 0xC0))
        return slide

    # 表紙
    cover = add_slide(str(title), first=True)
    sub = cover.shapes.add_textbox(margin, Inches(3.8), body_width, Inches(0.6))
    run = sub.text_frame.paragraphs[0].add_run()
    run.text = f"作成日：{datetime.now().strftime('%Y年%m月%d日')}　（法令・届出施設確認サポートAI）"
    _pptx_set_jp(run, 12, color=(0x66, 0x66, 0x66))

    for slide_title, blocks in _split_slides(_parse_blocks(md)):
        # 本文を「行」単位に平坦化してから上限行数で複数スライドに割る
        lines: list = []   # ("bullet"|"para"|"h3", text) / ("table", rows)
        for kind, payload in blocks:
            if kind == "heading":          # h3以下は小見出し行として扱う
                lines.append(("h3", payload[1]))
            elif kind == "bullets":
                lines += [("bullet", x) for x in payload]
            elif kind == "numbers":
                lines += [("bullet", f"{n}. {x}") for n, x in enumerate(payload, 1)]
            elif kind == "table":
                lines.append(("table", payload))
            else:
                lines += [("para", x) for x in str(payload).split("\n") if x.strip()]

        chunks: list = []
        cur: list = []
        for ln in lines:
            # 表は1枚に1つ。テキストは上限行数で分割する
            if ln[0] == "table" and cur or len(cur) >= _PPTX_MAX_LINES:
                chunks.append(cur)
                cur = []
            cur.append(ln)
            if ln[0] == "table":
                chunks.append(cur)
                cur = []
        if cur:
            chunks.append(cur)

        for ci, chunk in enumerate(chunks):
            slide = add_slide(slide_title if ci == 0 else f"{slide_title}（続き）")
            if chunk and chunk[0][0] == "table":
                rows = chunk[0][1]
                n_rows, n_cols = len(rows), max(len(r) for r in rows)
                shape = slide.shapes.add_table(
                    n_rows, n_cols, margin, Inches(1.9), body_width,
                    Inches(min(0.45 * n_rows, 5.0)))
                for ri, row in enumerate(rows):
                    for ci2 in range(n_cols):
                        cell = shape.table.cell(ri, ci2)
                        p = cell.text_frame.paragraphs[0]
                        _pptx_add_runs(p, row[ci2] if ci2 < len(row) else "",
                                       size_pt=12 if ri else 13)
                continue
            box = slide.shapes.add_textbox(margin, Inches(1.9), body_width, Inches(5.0))
            tf = box.text_frame
            tf.word_wrap = True
            first_p = True
            for kind2, text in chunk:
                p = tf.paragraphs[0] if first_p else tf.add_paragraph()
                first_p = False
                if kind2 == "h3":
                    _pptx_add_runs(p, re.sub(r"\*\*", "", text), size_pt=18)
                    for r in p.runs:
                        _pptx_set_jp(r, 18, bold=True, color=(0x15, 0x65, 0xC0))
                elif kind2 == "bullet":
                    _pptx_add_runs(p, "・" + text, size_pt=15)
                else:
                    _pptx_add_runs(p, text, size_pt=15)
                p.space_after = Pt(8)

    # 留意事項
    last = add_slide("留意事項")
    box = last.shapes.add_textbox(margin, Inches(1.9), body_width, Inches(3.0))
    box.text_frame.word_wrap = True
    _pptx_add_runs(box.text_frame.paragraphs[0], "⚠️ " + _DOC_NOTE, size_pt=13)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── PDF（PyMuPDF Story ＋ 日本語フォント埋め込み） ──────────

def build_pdf_file(title: str, md: str) -> bytes:
    import pymupdf as fitz

    font_path = next((p for p in _JP_FONT_CANDIDATES if os.path.exists(p)), None)
    if not font_path:
        raise RuntimeError("日本語フォントが見つかりません（PDF出力にはWindows標準の日本語フォントが必要です）")

    body = _blocks_to_html(_parse_blocks(md))
    html_doc = (
        f"<h1>{_inline_html(title)}</h1>"
        f'<p class="meta">作成日：{datetime.now().strftime("%Y年%m月%d日")}'
        f"　（法令・届出施設確認サポートAI）</p>"
        f"{body}"
        f'<p class="note">⚠️ {_DOC_NOTE}</p>'
    )
    css = (
        "@font-face { font-family: jp; src: url(%s); }"
        "* { font-family: jp; font-size: 10pt; line-height: 1.6; }"
        "h1 { font-size: 16pt; color: #1565C0; }"
        "h2 { font-size: 12pt; color: #1565C0; margin-top: 14pt; }"
        "h3 { font-size: 11pt; }"
        "p.meta, p.note { font-size: 8pt; color: #555; }"
        "th, td { border: 0.5pt solid #999; padding: 3pt 6pt; font-size: 9pt; }"
        % os.path.basename(font_path)
    )
    story = fitz.Story(html=html_doc, user_css=css,
                       archive=fitz.Archive(os.path.dirname(font_path)))
    buf = io.BytesIO()
    writer = fitz.DocumentWriter(buf)
    mediabox = fitz.paper_rect("a4")
    where = mediabox + (40, 44, -40, -44)
    more = 1
    while more:
        dev = writer.begin_page(mediabox)
        more, _ = story.place(where)
        story.draw(dev)
        writer.end_page()
    writer.close()

    # フォントファイル全体（游ゴシックで10MB超）が埋め込まれるため、
    # 使用グリフのみに間引いてサイズを数十KB台に抑える（要 fontTools）
    doc = fitz.open("pdf", buf.getvalue())
    try:
        doc.subset_fonts()
    except Exception:
        logger.warning("PDFフォントのサブセット化に失敗（サイズが大きくなります）", exc_info=True)
    return doc.tobytes(deflate=True, garbage=3)


def build_file(fmt: str, title: str, md: str) -> tuple:
    """(ファイル名, bytes, MIMEタイプ) を返す。fmt: html / docx / pdf"""
    if fmt == "html":
        return safe_filename(title, "html"), build_html_file(title, md), MIME["html"]
    if fmt == "docx":
        return safe_filename(title, "docx"), build_docx_file(title, md), MIME["docx"]
    if fmt == "pdf":
        return safe_filename(title, "pdf"), build_pdf_file(title, md), MIME["pdf"]
    if fmt == "pptx":
        return safe_filename(title, "pptx"), build_pptx_file(title, md), MIME["pptx"]
    raise ValueError(f"未対応の出力形式: {fmt}")
