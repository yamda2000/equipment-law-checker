"""アップロード資料からの設備情報抽出

PDF / Word / Excel / PowerPoint / テキストの資料からテキストを取り出し、
LLM でヒアリング11項目に該当する情報だけを構造化抽出する。
資料の全文は会話履歴に載せず、抽出結果（11項目の値と根拠）のみを後段に渡す。
"""

import io
import os
import re
import logging

from pydantic import BaseModel, Field
from langchain_core.messages import SystemMessage, HumanMessage

from backend.prompts import DOC_EXTRACTION_SYSTEM

logger = logging.getLogger(__name__)

# 1回のLLM抽出に渡すチャンクの文字数。長い資料はこのサイズで分割し、
# パートごとに抽出してから項目を統合する（末尾切り捨てによる後半の取りこぼしを防ぐ）。
# モデルのコンテキスト窓に安全に収まる範囲にする（gpt-4o 等で 45,000 字が目安）。
CHUNK_CHARS = int(os.getenv("DOC_CHUNK_CHARS", "45000"))

# 全資料を通しての絶対上限（暴走的なコスト・処理時間を防ぐ安全弁）。
# 通常はチャンク分割で全文を読むため切り捨ては起きない。これを超える巨大な
# アップロードのみ、超過分を末尾から切り捨てて警告する。
# 分割抽出のLLM呼び出し回数の上限はおおよそ MAX_TEXT_CHARS / CHUNK_CHARS。
MAX_TEXT_CHARS = int(os.getenv("DOC_MAX_TEXT_CHARS", "480000"))


# ─── ファイル種別ごとのテキスト抽出 ──────────────────────────────
def _extract_pdf(data: bytes) -> str:
    """PDF からテキストを抽出する（PyMuPDF）。カタログ等、埋め込みフォントの
    ToUnicode CMap が不完全な PDF でも文字化けしにくく高速。"""
    import fitz  # PyMuPDF
    parts = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for i, page in enumerate(doc, 1):
            text = page.get_text().strip()
            if text:
                parts.append(f"【{i}ページ】\n{text}")
    return "\n".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            parts.append(f"【シート: {ws.title}】")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"【スライド{i}】")
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_txt(data: bytes) -> str:
    for enc in ("utf-8-sig", "cp932"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_text_from_file(filename: str, data: bytes) -> str:
    """ファイル名の拡張子に応じてテキストを抽出する。
    抽出できない場合（スキャン画像PDF・破損ファイル等）は空文字を返す。"""
    name = filename.lower()
    try:
        if name.endswith(".pdf"):
            return _extract_pdf(data)
        if name.endswith(".docx"):
            return _extract_docx(data)
        if name.endswith((".xlsx", ".xlsm")):
            return _extract_xlsx(data)
        if name.endswith(".pptx"):
            return _extract_pptx(data)
        if name.endswith(".txt"):
            return _extract_txt(data)
    except Exception:
        logger.exception("資料テキスト抽出に失敗: %s", filename)
        return ""
    return ""


# ─── LLM 構造化抽出 ───────────────────────────────────────────────
class ExtractedField(BaseModel):
    value: str = Field(
        default="",
        description="資料に記載があった場合のみ、その内容を簡潔に記載（1〜2行）。記載がなければ空文字",
    )
    evidence: str = Field(
        default="",
        description="根拠となる資料名・記載箇所（例：仕様書.pdf 2ページ「排気」欄）。value が空なら空文字",
    )


class DocExtraction(BaseModel):
    equipment_type:     ExtractedField = Field(default_factory=ExtractedField, description="設備の種類")
    installation_place: ExtractedField = Field(default_factory=ExtractedField, description="設置場所（建屋・階・部屋名）")
    operation_purpose:  ExtractedField = Field(default_factory=ExtractedField, description="設備の用途・目的")
    scheduled_date:     ExtractedField = Field(default_factory=ExtractedField, description="稼働開始予定日")
    chemicals:          ExtractedField = Field(default_factory=ExtractedField, description="薬品・溶剤・ガス・燃料の使用有無と種類・使用量・貯蔵量。チラー等に内蔵の冷媒ガス（フロン・R410A等）・圧縮ガスも必ずここに記載（例：冷媒R410A 充填量2.5kg）")
    fire_exhaust:       ExtractedField = Field(default_factory=ExtractedField, description="火気・熱源・排気・粉じんの発生有無")
    wastewater:         ExtractedField = Field(default_factory=ExtractedField, description="排水・廃液・廃棄物の発生有無")
    noise_vibration:    ExtractedField = Field(default_factory=ExtractedField, description="騒音・振動の発生有無")
    radiation:          ExtractedField = Field(default_factory=ExtractedField, description="放射線・X線発生装置への該当有無")
    construction:       ExtractedField = Field(default_factory=ExtractedField, description="建屋改修・電気工事（受電容量増加・自家発電機・蓄電池含む）・配管工事の有無")
    additional_info:    ExtractedField = Field(default_factory=ExtractedField, description="上記以外の設備導入関連情報（仕様・重量・搬入経路・メーカー名など。冷媒・ガス類は chemicals に記載する）")


# 冷媒・ガス類の検出パターン（フロン排出抑制法・高圧ガス保安法の判定材料）
_REFRIGERANT_PAT = re.compile(r"冷媒|フロン|R\d{2,3}[A-Za-z]?\b|HFC|高圧ガス")


def _relocate_refrigerant(data: dict) -> dict:
    """LLMが冷媒・ガス類を additional_info に振り分けた場合、chemicals に補記する。
    （プロンプト指示だけでは安定しないため、確定的に後処理する。
    additional_info 側は他の情報を含み得るためそのまま残し、確認画面で担当者が整理する）"""
    add = data.get("additional_info") or {}
    chem = data.get("chemicals") or {}
    add_val = add.get("value", "")
    if _REFRIGERANT_PAT.search(add_val) and not _REFRIGERANT_PAT.search(chem.get("value", "")):
        chem["value"] = (chem["value"] + "／" if chem.get("value") else "") + add_val
        if not chem.get("evidence"):
            chem["evidence"] = add.get("evidence", "")
        data["chemicals"] = chem
    return data


def _fit_doc_texts(doc_texts: list, cap: int) -> tuple[list, list]:
    """各資料に文字数上限を配分し、(採用リスト, 切り捨て情報) を返す。

    従来の「連結後に先頭 cap 字」方式では、1件目の大きな資料だけで上限を
    使い切ると2件目以降が1文字もLLMへ渡らず、無通知の確認漏れになる。
    短い資料の未使用分を長い資料へ再配分しつつ、全資料に必ず枠を与える。
    """
    n = len(doc_texts)
    if n == 0:
        return [], []
    remaining = cap
    # 長さ昇順に配分すると、短い資料の余剰が自動的に長い資料へ回る
    order = sorted(range(n), key=lambda i: len(doc_texts[i][1]))
    alloc = [0] * n
    left = n
    for i in order:
        share = remaining // left
        alloc[i] = min(len(doc_texts[i][1]), share)
        remaining -= alloc[i]
        left -= 1
    fitted, truncated = [], []
    for i, (name, text) in enumerate(doc_texts):
        fitted.append((name, text[:alloc[i]]))
        if len(text) > alloc[i]:
            truncated.append(
                {"name": name, "used": alloc[i], "total": len(text)}
            )
    return fitted, truncated


def _split_text(text: str, size: int) -> list[str]:
    """text を size 文字以下の断片に分割する。できるだけ改行位置で区切り、
    1行が size を超える場合のみ途中で強制的に切る（ページ・行の途中割れを最小化）。"""
    if size <= 0 or len(text) <= size:
        return [text] if text else []
    pieces: list[str] = []
    i, n = 0, len(text)
    while i < n:
        if n - i <= size:
            pieces.append(text[i:])
            break
        cut = text.rfind("\n", i, i + size)
        if cut <= i:
            cut = i + size  # 改行が無ければ強制カット
        pieces.append(text[i:cut])
        i = cut
        while i < n and text[i] == "\n":  # 断片先頭の改行を食う
            i += 1
    return pieces


def _build_chunks(fitted: list, chunk_chars: int) -> list[str]:
    """採用テキストを、資料名ラベル付きで chunk_chars 以下のチャンクにまとめる。
    小さい資料は詰め合わせ、大きい資料は改行境界で分割して複数チャンクに渡す。"""
    chunks: list[str] = []
    cur: list[str] = []
    cur_len = 0
    for name, text in fitted:
        if not text.strip():
            continue
        header = f"===== 資料: {name} =====\n"
        budget = max(1000, chunk_chars - len(header))
        for piece in _split_text(text, budget):
            block = header + piece
            if cur and cur_len + len(block) > chunk_chars:
                chunks.append("\n\n".join(cur))
                cur, cur_len = [], 0
            cur.append(block)
            cur_len += len(block) + 2
    if cur:
        chunks.append("\n\n".join(cur))
    return chunks


def _norm(s: str) -> str:
    """統合時の重複判定用の正規化（空白除去・小文字化）。"""
    return re.sub(r"\s+", "", s or "").lower()


def _merge_extractions(partials: list[dict]) -> dict:
    """複数パートの抽出結果を項目ごとに統合する。
    同一・内包関係の値は重複としてまとめ、異なる値は「／」で併記する
    （後半パートの情報を落とさないことを優先し、確認画面で担当者が整理する想定）。"""
    merged: dict = {}
    for field in DocExtraction.model_fields:
        values: list[str] = []
        evidences: list[str] = []
        for p in partials:
            f = p.get(field) or {}
            v = (f.get("value") or "").strip()
            if not v:
                continue
            nv = _norm(v)
            dup = False
            for idx, existing in enumerate(values):
                ne = _norm(existing)
                if nv == ne or nv in ne:
                    dup = True            # 既出と同一・既出に内包される
                    break
                if ne in nv:              # 新しい値が既出を内包 → 置き換える
                    values[idx] = v
                    dup = True
                    break
            if not dup:
                values.append(v)
            e = (f.get("evidence") or "").strip()
            if e and e not in evidences:
                evidences.append(e)
        merged[field] = {
            "value": "／".join(values)[:1000],
            "evidence": "；".join(evidences)[:500],
        }
    return merged


def extract_equipment_info(
    doc_texts: list, config: dict | None = None
) -> tuple[dict, list]:
    """[(ファイル名, 抽出テキスト), ...] からヒアリング11項目を構造化抽出する。

    長い資料は CHUNK_CHARS 単位に分割して各パートから抽出し、項目ごとに統合する
    （従来の末尾切り捨てによる後半情報の取りこぼしを防ぐ）。全資料の合計が
    MAX_TEXT_CHARS を超える場合のみ、安全弁として超過分を切り捨てて報告する。

    config: LLM 呼び出しに渡す RunnableConfig（Langfuse の callbacks 等）。
    返り値: ({フィールド名: {"value": str, "evidence": str}},
             切り捨て情報 [{"name", "used", "total"}, ...])
    抽出に失敗した場合は例外を送出する（呼び出し側でフォールバック）。
    """
    from backend.workflow import _llm  # 循環importを避けるため遅延import

    fitted, truncated = _fit_doc_texts(doc_texts, MAX_TEXT_CHARS)
    if truncated:
        logger.warning(
            "資料テキストが絶対上限 %d 字を超過。切り捨て: %s",
            MAX_TEXT_CHARS,
            [(t["name"], f'{t["used"]}/{t["total"]}字') for t in truncated],
        )

    llm = _llm().with_structured_output(DocExtraction)

    def _extract(chunk_text: str) -> dict:
        result: DocExtraction = llm.invoke([
            SystemMessage(DOC_EXTRACTION_SYSTEM),
            HumanMessage(f"以下の資料から設備情報を抽出してください。\n\n{chunk_text}"),
        ], config=config or {})
        return {name: field.model_dump() for name, field in result}

    chunks = _build_chunks(fitted, CHUNK_CHARS)
    if not chunks:
        # 全資料が空（抽出テキストなし）：空の結果を返す
        merged = {name: field.model_dump() for name, field in DocExtraction()}
    elif len(chunks) == 1:
        merged = _extract(chunks[0])
    else:
        # 資料が長い：分割抽出 → 項目統合（後半の取りこぼしを防ぐ）
        logger.info("資料を %d チャンクに分割して抽出します（合計約%d字）",
                    len(chunks), sum(len(t) for _n, t in fitted))
        merged = _merge_extractions([_extract(c) for c in chunks])

    return _relocate_refrigerant(merged), truncated
